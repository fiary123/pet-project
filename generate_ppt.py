import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ================= 配置区 =================
OUTPUT_FILE = '大学信息技术基础_期末汇报.pptx'
TEMPLATE_FILE = r"C:\Users\CX10\Desktop\大学信息技术期末模板.pptx"  # 确保你的模板文件叫这个名字

# 学生信息 (来自 Word)
INFO = {
    "title": "大学信息技术基础学习报告",
    "subtitle": "2025日语1班  李雪萌\n2025104210136",
    "date": "2025年12月"
}

# ================= 核心工具函数 =================

def set_font(run, size, is_bold=True, color=None):
    """设置字体格式：默认微软雅黑，支持加粗和颜色"""
    font = run.font
    font.name = '微软雅黑'  # 演示文稿常用字体
    font.size = Pt(size)
    font.bold = is_bold
    if color:
        font.color.rgb = color
    else:
        font.color.rgb = RGBColor(0, 0, 0) # 默认纯黑

def format_title(shape, text, level):
    """
    根据用户要求严格设置标题格式
    Level 1: 黑体 32pt 左顶格
    Level 2: 黑体 24pt 左顶格
    Level 3: 黑体 20pt 左顶格
    """
    shape.text = text
    tf = shape.text_frame
    
    # 确定字号
    if level == 1:
        font_size = 32
    elif level == 2:
        font_size = 24
    else:
        font_size = 20
        
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT # 左顶格
        for run in p.runs:
            set_font(run, font_size, is_bold=True) # 黑体(加粗)

def get_safe_layout(prs, index):
    """
    安全获取版式：如果指定的 index 不存在，就退回到 index 1 (通常是正文页)
    """
    try:
        return prs.slide_layouts[index]
    except IndexError:
        print(f"提示：模板中找不到索引为 {index} 的版式，自动退回到索引 1。")
        return prs.slide_layouts[1]

def add_content_slide(prs, title, level, bullets, img_path=None):
    """
    添加一张正文页
    """
    # 使用安全版式获取（尝试获取索引1，即内容页）
    slide = prs.slides.add_slide(get_safe_layout(prs, 1))
    
    # 1. 设置标题
    # 有些模板的 layout[1] 可能没有 title，做个判断
    if slide.shapes.title:
        format_title(slide.shapes.title, title, level)
    
    # 2. 处理内容区
    # 获取正文占位符 (通常是 placeholders[1])
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        
        # 如果有图片且存在
        has_image = img_path and os.path.exists(img_path)
        
        if has_image:
            # 调整文本框宽度，左侧留给文字 (约占50%)
            body_shape.width = Inches(5.0) 
            body_shape.height = Inches(5.5)
            body_shape.top = Inches(1.8)
            body_shape.left = Inches(0.5)
            
            # 插入图片在右侧
            slide.shapes.add_picture(img_path, left=Inches(5.8), top=Inches(2.0), height=Inches(4.0))
        else:
            # 无图模式，文本稍微居中或宽一点
            body_shape.top = Inches(1.8)
            body_shape.width = Inches(9.0)

        # 填充文本
        tf = body_shape.text_frame
        tf.clear() # 清除默认格式
        
        for item in bullets:
            p = tf.add_paragraph()
            # 判断是否是子层级
            if isinstance(item, dict): 
                 text = item['text']
                 lvl = item.get('lvl', 0)
            else:
                 text = item
                 lvl = 0
                 
            p.text = text
            p.level = lvl
            p.space_before = Pt(12) # 段间距
            p.space_after = Pt(6)
            
            # 正文字体设置 (20pt 比较清晰)
            for run in p.runs:
                set_font(run, 20 if lvl==0 else 18, is_bold=False)
    else:
        print(f"警告：页面 '{title}' 的版式似乎没有正文占位符，跳过文字填充。")

    return slide

# ================= 主程序逻辑 =================

def create_full_ppt():
    # 1. 加载模板
    if os.path.exists(TEMPLATE_FILE):
        prs = Presentation(TEMPLATE_FILE)
        print(f"成功加载模板: {TEMPLATE_FILE}")
        # 打印一下有多少个 layout，方便调试
        print(f"当前模板包含 {len(prs.slide_layouts)} 种版式。")
    else:
        print("未找到模板文件，正在创建空白演示文稿...")
        prs = Presentation()

    # --- Slide 1: 封面 ---
    # 使用 layout 0 (通常是封面)
    slide = prs.slides.add_slide(get_safe_layout(prs, 0))
    if slide.shapes.title:
        slide.shapes.title.text = INFO["title"]
    # 尝试找副标题占位符
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = INFO["subtitle"]

    # --- Slide 2: 目录 ---
    slide = prs.slides.add_slide(get_safe_layout(prs, 1))
    if slide.shapes.title:
        format_title(slide.shapes.title, "目 录", 1)
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "1. 引言\n2. 理论知识\n3. 实践经验\n4. 学习成果\n5. 总结与展望"
        for p in tf.paragraphs:
            p.space_before = Pt(20)
            for run in p.runs:
                set_font(run, 24, True)

    # --- Slide 3: 1 引言 ---
    content = [
        "背景：信息技术已超越工具范畴，成为重塑人类认知的基础设施。",
        "“人-机-物”三元融合正在加速实现。",
        "学习目标：",
        {'text': "不只是掌握软件操作", 'lvl': 1},
        {'text': "构建核心竞争力：信息素养 (Information Literacy)", 'lvl': 1},
        {'text': "培养计算思维 (Computational Thinking)", 'lvl': 1}
    ]
    add_content_slide(prs, "1 引言", 1, content)

    # --- Slide 4: 2.1.1 信息技术概念 ---
    content = [
        "定义演变：",
        {'text': "1958年：处理大量信息的技术与统计数学方法的结合。", 'lvl': 1},
        {'text': "现代：利用计算机系统采集、存储、处理、传输信息。", 'lvl': 1},
        "核心价值：从“数据处理”到“智慧赋能”。",
        "使能技术：重塑教育、医疗、金融等行业形态（如智能辅导系统）。"
    ]
    add_content_slide(prs, "2.1.1 信息技术的基本概念", 3, content)

    # --- Slide 5: 2.1.2 硬件体系 ---
    content = [
        "冯·诺依曼架构 (Von Neumann Architecture)：",
        {'text': "二进制逻辑", 'lvl': 1},
        {'text': "存储程序概念 (核心思想)", 'lvl': 1},
        "五大核心部件：",
        {'text': "运算器 (ALU) & 控制器 (CU) -> CPU", 'lvl': 1},
        {'text': "存储器 (Memory/Storage)", 'lvl': 1},
        {'text': "输入设备 & 输出设备", 'lvl': 1}
    ]
    add_content_slide(prs, "2.1.2 硬件系统：冯·诺依曼架构", 3, content, "von_neumann.png")

    # --- Slide 6: 2.1.3 物联网架构 ---
    content = [
        "IoT 三层技术架构：",
        "1. 感知层 (Perception Layer)：",
        {'text': "“皮肤与五官”，负责采集信息 (传感器, RFID)。", 'lvl': 1},
        "2. 网络层 (Network Layer)：",
        {'text': "负责数据传输 (5G, Wi-Fi, LPWAN)。", 'lvl': 1},
        "3. 应用层 (Application Layer)：",
        {'text': "结合行业需求 (智慧农业, 智慧城市)。", 'lvl': 1}
    ]
    add_content_slide(prs, "2.1.3 物联网的技术特征", 3, content, "iot_structure.png")

    # --- Slide 7: 物联网通信技术对比 (表格页) ---
    # 【修复重点】这里原来是 layout[5]，改为安全获取 layout[1]
    slide = prs.slides.add_slide(get_safe_layout(prs, 1)) 
    if slide.shapes.title:
        format_title(slide.shapes.title, "2.1.3 主流无线通信技术对比", 3)
    
    # 绘制表格
    rows, cols = 5, 4
    left = Inches(0.5); top = Inches(2.0); width = Inches(9.0); height = Inches(3.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    headers = ['通信技术', '覆盖范围', '功耗', '典型应用场景']
    data = [
        ['Zigbee', '短 (10-100m)', '极低', '智能家居（灯光、窗帘）'],
        ['Wi-Fi', '中短 (<100m)', '高', '视频监控、大数据传输'],
        ['NB-IoT', '广 (>10km)', '低', '智能抄表、井盖监测'],
        ['5G', '广域', '中高', '自动驾驶、远程手术']
    ]
    
    # 填充表格
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(50, 100, 200) # 蓝色表头
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_font(run, 18, True, RGBColor(255,255,255))
                
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx+1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_font(run, 16, False)

    # --- Slide 8: 2.2 云计算 ---
    content = [
        "核心定义：计算资源的“水与电”。",
        "六大特征 (NIST)：",
        {'text': "资源池化 (Resource Pooling)", 'lvl': 1},
        {'text': "快速弹性 (Rapid Elasticity) - 应对波峰波谷", 'lvl': 1},
        {'text': "按需自助服务", 'lvl': 1},
        "SPI 服务模式：",
        {'text': "IaaS (基础设施), PaaS (平台), SaaS (软件)", 'lvl': 1}
    ]
    add_content_slide(prs, "2.2 云计算：算力的资源化", 2, content)

    # --- Slide 9: 2.3 大数据 ---
    content = [
        "5V 特征：",
        {'text': "Volume (大量), Velocity (高速)", 'lvl': 1},
        {'text': "Variety (多样), Veracity (真实), Value (价值)", 'lvl': 1},
        "技术生态：",
        {'text': "Hadoop (离线存储与计算)", 'lvl': 1},
        {'text': "Spark (内存计算，实时分析)", 'lvl': 1}
    ]
    add_content_slide(prs, "2.3 大数据技术", 2, content)

    # --- Slide 10: 2.4 人工智能 ---
    content = [
        "技术突破 (2024-2025)：",
        {'text': "生成式AI (Generative AI) 崛起 (ChatGPT, Gemini)。", 'lvl': 1},
        {'text': "基于 Transformer 架构与自注意力机制。", 'lvl': 1},
        "行业应用：",
        {'text': "医疗：影像分析、新药研发。", 'lvl': 1},
        {'text': "金融：高频交易、风控。", 'lvl': 1}
    ]
    add_content_slide(prs, "2.4 人工智能", 2, content)

    # --- Slide 11: 3.1 Word 实践 (问题) ---
    content = [
        "项目：本科毕业论文排版（30页）",
        "遇到的问题：",
        {'text': "手动编号混乱：插入新章节后，后续编号需全部重调。", 'lvl': 1},
        {'text': "目录生成失败：Word无法识别手动设置的“大号字体”为标题。", 'lvl': 1},
        {'text': "页眉控制困难：无法实现“目录罗马数字，正文阿拉伯数字”的页码切换。", 'lvl': 1}
    ]
    add_content_slide(prs, "3.1 Word长文档：遇到的挑战", 2, content)

    # --- Slide 12: 3.1 Word 实践 (解决) ---
    content = [
        "核心解决方案：",
        "1. 样式 (Styles) 系统：",
        {'text': "修改“标题1”等内置样式，实现内容与格式分离。", 'lvl': 1},
        {'text': "定义多级列表，实现“第1章 1.1”自动编号。", 'lvl': 1},
        "2. 分节符 (Section Break)：",
        {'text': "断开节链接，独立控制目录页和正文页的页码。", 'lvl': 1},
        "3. 域 (Field) 应用：",
        {'text': "使用 StyleRef 域在页眉动态引用当前章标题。", 'lvl': 1}
    ]
    add_content_slide(prs, "3.1 Word长文档：工程化解决方案", 2, content)

    # --- Slide 13: 3.2 Excel 实践 (问题) ---
    content = [
        "项目：连锁超市销售数据分析",
        "遇到的问题：",
        {'text': "数据脏乱：VLOOKUP 匹配 #N/A，原因是肉眼不可见的尾部空格。", 'lvl': 1},
        {'text': "分析效率低：使用 SUMIFS 写长公式，容易出错且难以修改。", 'lvl': 1},
        {'text': "可视化误区：使用了 3D 饼图，分类过多导致难以阅读。", 'lvl': 1}
    ]
    add_content_slide(prs, "3.2 Excel数据分析：遇到的挑战", 2, content)

    # --- Slide 14: 3.2 Excel 实践 (解决) ---
    content = [
        "核心解决方案：",
        "1. ETL 数据清洗：",
        {'text': "TRIM() 去空格，分列功能修复数字格式。", 'lvl': 1},
        "2. 数据透视表 (Pivot Table)：",
        {'text': "拖拽生成多维报表，配合切片器实现动态交互。", 'lvl': 1},
        "3. 科学可视化：",
        {'text': "使用组合图 (柱状+折线) 展示销售额与增长率。", 'lvl': 1},
        {'text': "遵循“少即是多”原则，去除多余装饰。", 'lvl': 1}
    ]
    add_content_slide(prs, "3.2 Excel数据分析：商业智能应用", 2, content, "excel_chart.png")

    # --- Slide 15: 3.3 PPT 实践 ---
    content = [
        "设计原则：",
        {'text': "一页一观点，拒绝大段文字堆砌。", 'lvl': 1},
        {'text': "CRAP原则：对比、重复、对齐、亲密性。", 'lvl': 1},
        "技术应用：",
        {'text': "幻灯片母版 (Slide Master) 统一字体与Logo。", 'lvl': 1},
        {'text': "SmartArt 将文本列表转化为逻辑图示。", 'lvl': 1},
        {'text': "备注栏：将详细讲稿移至备注，屏幕只留核心观点。", 'lvl': 1}
    ]
    add_content_slide(prs, "3.3 演示文稿：视觉传达与逻辑", 2, content)

    # --- Slide 16: 4 学习成果 (Word/Excel) ---
    content = [
        "Word 工程化思维：",
        {'text': "从“格式刷”进阶到“样式库”。", 'lvl': 1},
        {'text': "掌握长文档的自动化管理 (目录、题注、交叉引用)。", 'lvl': 1},
        "Excel 数据逻辑：",
        {'text': "深刻理解“Garbage In, Garbage Out”。", 'lvl': 1},
        {'text': "掌握从数据清洗到动态仪表盘的全流程。", 'lvl': 1}
    ]
    add_content_slide(prs, "4 学习成果：办公软件进阶", 1, content)

    # --- Slide 17: 4 学习成果 (综合) ---
    content = [
        "PPT 视觉叙事：",
        {'text': "从“文档搬运工”转变为“视觉设计师”。", 'lvl': 1},
        {'text': "掌握母版设计与 Morph 平滑切换动画。", 'lvl': 1},
        "信息素养提升：",
        {'text': "学会利用 AI 辅助学习 (查报错、读文档)。", 'lvl': 1},
        {'text': "建立解决复杂问题的计算思维。", 'lvl': 1}
    ]
    add_content_slide(prs, "4 学习成果：思维与素养", 1, content)

    # --- Slide 18: 5 总结与展望 ---
    content = [
        "从“工具”到“思维”：",
        {'text': "Excel 函数嵌套蕴含算法思想，Word 样式体现结构化编程。", 'lvl': 1},
        "技术理性与人文关怀：",
        {'text': "关注算法偏见、隐私保护与数字鸿沟。", 'lvl': 1},
        "终身学习：",
        {'text': "技术迭代极快 (如 Sora, DeepSeek)。", 'lvl': 1},
        {'text': "核心能力是“如何学习新工具”的自驱力。", 'lvl': 1}
    ]
    add_content_slide(prs, "5 总结与展望", 1, content)

    # --- Slide 19: 结束页 ---
    slide = prs.slides.add_slide(get_safe_layout(prs, 0)) # 结束页也用 layout 0 (类似封面)
    if slide.shapes.title:
        slide.shapes.title.text = "谢谢观看！"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "请老师批评指正"

    # 保存
    prs.save(OUTPUT_FILE)
    print(f"PPT生成完毕，共 {len(prs.slides)} 页。文件名为: {OUTPUT_FILE}")

if __name__ == "__main__":
    create_full_ppt()